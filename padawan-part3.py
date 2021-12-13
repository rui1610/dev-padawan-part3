from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]

###################################################################################
# Add a slide
###################################################################################
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
#################################################################
# Add a text to the slide
#################################################################
title.text = "Hello, my very dear Padawan!"
subtitle.text = "I can feel The Force is strong inside of you!"
left =  Inches(1)
top = Inches(6)
width = Inches(5)
height = Inches(1)
textBox = slide.shapes.add_textbox(left, top, width, height)
#################################################################
# Add another text to the slide
#################################################################
tf = textBox.text_frame
tf.text = "But never forget:"
tf.paragraphs[0].font.size = Pt(20)
tf.paragraphs[0].font.bold = True
#################################################################
# ... and that text has a text with a link, too
#################################################################
p = tf.add_paragraph()
run = p.add_run()
run.text = "Fear leads to the dark side of The Force!"
run.font.size = Pt(20)
run.hyperlink.address = "https://www.youtube.com/watch?v=kFnFr-DOPf8"


###################################################################################
# Add another slide here, if you want
###################################################################################



filename = 'The force inside of you is strong.pptx'
prs.save(filename)


